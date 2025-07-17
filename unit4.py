from pptx import Presentation
from pptx.util import Inches, Pt

# Create presentation and define styles
prs = Presentation()
blank_layout = prs.slide_layouts[6]  # blank for consistent dark-theme style

def add_slide(title, content_lines):
    slide = prs.slides.add_slide(blank_layout)
    # Title
    tx = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    p = tx.text_frame
    p.text = title
    p.paragraphs[0].font.size = Pt(32)
    p.paragraphs[0].font.bold = True

    # Body text
    text_top = Inches(1.5)
    for line in content_lines:
        tb = slide.shapes.add_textbox(Inches(1), text_top, Inches(8), Inches(1))
        pf = tb.text_frame.add_paragraph()
        pf.text = line
        pf.font.size = Pt(20)
        text_top += Inches(0.6)

# Slide 1
add_slide("Unit 4 – Coding & Circuits", ["From Python to Qubits"])

# Slide 2
add_slide("Why Coding for Quantum?",
    ["• We simulate quantum circuits with Python",
     "• It's like rehearsing before the real experiment",
     "• Outcome: by the end, you'll write & run your own circuits"])

# Slide 3
add_slide("Python Crash Course (Prework)",
    ["Watch before class: intro to variables, lists, functions",
     "We'll only use these basics today"])

# Slide 4
add_slide("Hello Qiskit",
    ["```python",
     "from qiskit import QuantumCircuit",
     "qc = QuantumCircuit(1)",
     "qc.h(0)",
     "qc.measure_all()",
     "print(qc.draw())",
     "```"])

# Slide 5
add_slide("What Is a Classical Circuit?",
    ["• Bits → Gates → Output",
     "• Gates act on 0 or 1",
     "• Many are irreversible (e.g. AND)"])

# Slide 6
add_slide("What Is a Quantum Circuit?",
    ["• Qubits → Quantum gates → Measurement",
     "• Gates here are *reversible matrices*",
     "• Qubits can be in superposition until measured"])

# Slide 7
add_slide("Classical vs Quantum Gates",
    ["Classical: AND (❌ irreversible)",
     "Quantum: X/NOT, CNOT/XOR (✅ reversible)"])

# Slide 8
add_slide("Anatomy of a Quantum Script",
    ["1. Import libs",
     "2. Define circuit",
     "3. Add gates",
     "4. Measure",
     "5. Run & collect results"])

# Slide 9 (Quiz)
add_slide("In‑Class Quiz",
    ["1. What does `qc.h(0)` do? (A) Apply H on qubit 0 ✅",
     "2. Which is quantum‑specific? (B) Measurement ✅"])

# Slide 10
add_slide("Summary & Next Steps",
    ["• Python lets us define & simulate quantum circuits",
     "• Quantum gates are reversible matrix ops",
     "• Measurement gives classical output (0 or 1)",
     "→ Up next: Unit 5 – Quantum Algorithms"])

# Save the generated pptx
prs.save("Unit4_Coding_Circuits.pptx")
print("✅ Presentation saved to Unit4_Coding_Circuits.pptx")
